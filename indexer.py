import os
import sys
import sqlite3
import time
import datetime
from tqdm import tqdm
import traceback
import logging

# --- File Processing Libs ---
try:
    import docx  # python-docx
except ImportError:
    print("Warning: python-docx not found. Install it (`pip install python-docx`) to process .docx files.", file=sys.stderr)
    docx = None
try:
    import openpyxl
except ImportError:
    print("Warning: openpyxl not found. Install it (`pip install openpyxl`) to process .xlsx files.", file=sys.stderr)
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    print("Warning: python-pptx not found. Install it (`pip install python-pptx`) to process .pptx files.", file=sys.stderr)
    Presentation = None
try:
    import fitz  # PyMuPDF
except ImportError:
    print("Warning: PyMuPDF not found. Install it (`pip install PyMuPDF`) to process .pdf files.", file=sys.stderr)
    fitz = None
try:
    from PIL import Image, UnidentifiedImageError
except ImportError:
    print("Warning: Pillow not found. Install it (`pip install Pillow`) for image processing.", file=sys.stderr)
    Image = None
try:
    import pytesseract # Requires Tesseract OCR installed
except ImportError:
    print("Warning: pytesseract not found. Install it (`pip install pytesseract`) for OCR.", file=sys.stderr)
    pytesseract = None

# --- Text Processing Libs ---
try:
    import nltk
    from nltk.corpus import stopwords
    from nltk.tokenize import word_tokenize, sent_tokenize
    from nltk.probability import FreqDist
except ImportError:
    print("Error: NLTK not found. Install it (`pip install nltk`) for summarization/keywords.", file=sys.stderr)
    sys.exit(1)
try:
    # --- Explicitly add VENV path to NLTK data path ---
    venv_nltk_path = os.path.join(os.path.dirname(__file__), '.venv', 'nltk_data')
    if os.path.isdir(venv_nltk_path):
        if venv_nltk_path not in nltk.data.path:
            nltk.data.path.append(venv_nltk_path)
            print(f"Added VENV NLTK path: {venv_nltk_path}")
    
    from sumy.parsers.plaintext import PlaintextParser
    from sumy.nlp.tokenizers import Tokenizer as SumyTokenizer
    from sumy.summarizers.lsa import LsaSummarizer
    from sumy.nlp.stemmers import Stemmer
    from sumy.utils import get_stop_words
except ImportError:
    print("Warning: sumy not found. Install it (`pip install sumy`) for summarization.", file=sys.stderr)
    PlaintextParser = None # Allow script to run without sumy, skipping summarization

# --- Configuration --
DATABASE_NAME = 'file_index.db'
SUMMARY_SENTENCE_COUNT = 3  # Number of sentences for the summary
MAX_KEYWORDS = 10 # Max number of keywords to extract
LANGUAGE = "english" # For summarizer and keyword extraction
LOG_FILE = 'indexing_errors.log'
COMMIT_INTERVAL = 100 # Commit to DB every N files

# --- Logging Setup ---
logging.basicConfig(filename=LOG_FILE, level=logging.WARNING,
                    format='%(asctime)s - %(levelname)s - %(message)s')

# --- NLTK Data Check ---
def check_nltk_data():
    try:
        nltk.data.find('tokenizers/punkt')
        nltk.data.find('corpora/stopwords')
        return True
    except LookupError:
        print("\nError: Required NLTK data (punkt, stopwords) not found.", file=sys.stderr)
        print("Please download them by running:", file=sys.stderr)
        print("python -m nltk.downloader punkt stopwords", file=sys.stderr)
        return False
    except Exception as e:
        print(f"\nError checking NLTK data: {e}", file=sys.stderr)
        return False

# --- Helper Functions ---

def setup_database(db_name=DATABASE_NAME):
    """Creates the SQLite database and the main table if they don't exist."""
    try:
        conn = sqlite3.connect(db_name, timeout=30)
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                path TEXT UNIQUE NOT NULL,
                filename TEXT NOT NULL,
                extension TEXT,
                size_bytes INTEGER,
                last_modified REAL, -- Store as Unix timestamp
                category_year INTEGER,
                category_type TEXT,
                category_event TEXT DEFAULT 'Unknown', -- Placeholder with default
                category_meeting TEXT DEFAULT 'Unknown', -- Placeholder with default
                summary TEXT,
                keywords TEXT, -- Store as comma-separated string
                processing_status TEXT DEFAULT 'Pending', -- Pending, Success, Failed
                processing_error TEXT -- Store error message if processing failed
            )
        ''')
        # Add indexes for faster searching
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_path ON files (path)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_filename ON files (filename)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_type ON files (category_type)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_year ON files (category_year)')
        cursor.execute('CREATE INDEX IF NOT EXISTS idx_status ON files (processing_status)')
        conn.commit()
        return conn, cursor
    except sqlite3.Error as e:
        print(f"FATAL: Database setup failed: {e}", file=sys.stderr)
        logging.critical(f"Database setup failed: {e}")
        sys.exit(1)


def get_file_type(extension):
    """Basic categorization based on file extension."""
    ext = extension.lower() if extension else ''
    if ext in ['.txt', '.log', '.md', '.csv', '.rtf']:
        return 'Text'
    elif docx and ext in ['.doc', '.docx']:
        return 'Word Document'
    elif openpyxl and ext in ['.xls', '.xlsx']:
        return 'Excel Spreadsheet'
    elif Presentation and ext in ['.ppt', '.pptx']:
        return 'PowerPoint Presentation'
    elif fitz and ext == '.pdf':
        return 'PDF Document'
    elif Image and ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        return 'Image'
    elif ext in ['.zip', '.rar', '.tar', '.gz', '.7z']:
        return 'Archive'
    # Add more specific types if needed
    elif ext in ['.py', '.js', '.java', '.c', '.cpp', '.h', '.cs', '.html', '.css', '.sh']:
        return 'Code'
    elif ext in ['.mp3', '.wav', '.aac', '.flac', '.ogg']:
        return 'Audio'
    elif ext in ['.mp4', '.avi', '.mkv', '.mov', '.wmv']:
        return 'Video'
    else:
        return 'Other'

# --- Text Extraction Functions ---
# Each function now returns text or None, logging errors internally

def extract_text_docx(file_path):
    if not docx: return None
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        logging.warning(f"Failed to extract text from DOCX '{file_path}': {e}")
        return None

def extract_text_xlsx(file_path):
    if not openpyxl: return None
    text_content = []
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in workbook:
            for row in sheet.iter_rows():
                row_text = [str(cell.value) for cell in row if cell.value is not None]
                if row_text:
                    text_content.append(" ".join(row_text))
        return "\n".join(text_content) if text_content else None
    except Exception as e:
        logging.warning(f"Failed to extract text from XLSX '{file_path}': {e}")
        return None

def extract_text_pptx(file_path):
    if not Presentation: return None
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content) if text_content else None
    except Exception as e:
        logging.warning(f"Failed to extract text from PPTX '{file_path}': {e}")
        return None

def extract_text_pdf(file_path):
    if not fitz: return None
    text_content = []
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text_content.append(page.get_text())
        doc.close()
        return "\n".join(text_content) if text_content else None
    except Exception as e:
        logging.warning(f"Failed to extract text from PDF '{file_path}': {e}")
        return None

def extract_text_image(file_path):
    if not Image or not pytesseract: return None
    try:
        # Check if tesseract executable is set, otherwise it checks PATH
        # You might need to set this explicitly:
        # pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract' # Example for Linux
        return pytesseract.image_to_string(Image.open(file_path), timeout=30) # Add timeout
    except (pytesseract.TesseractNotFoundError, FileNotFoundError):
        logging.error("Tesseract OCR engine not found or not in PATH. Cannot process images.")
        # Disable further attempts after first error? Maybe add a flag.
        return None
    except UnidentifiedImageError:
        logging.warning(f"Cannot identify image file: '{file_path}'")
        return None
    except Exception as e:
        logging.warning(f"Failed OCR on image '{file_path}': {e}")
        return None

def extract_text_plain(file_path):
    encodings_to_try = ['utf-8', 'latin-1', 'cp1252']
    for enc in encodings_to_try:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logging.warning(f"Failed to read plain text file '{file_path}' with encoding {enc}: {e}")
            return None # Give up if any read error occurs after trying encodings
    logging.warning(f"Could not decode plain text file '{file_path}' with tried encodings.")
    return None # If all encodings fail

# --- Text Processing ---
stop_words_set = set(stopwords.words(LANGUAGE))

def summarize_text(text, sentences_count=SUMMARY_SENTENCE_COUNT):
    if not text or not isinstance(text, str) or len(text.strip()) < 20: # Min length for meaningful summary
        return ""
    if not PlaintextParser: # Check if sumy is available
         logging.warning("Sumy library not available, skipping summarization.")
         return ""
    try:
        parser = PlaintextParser.from_string(text, SumyTokenizer(LANGUAGE))
        stemmer = Stemmer(LANGUAGE)
        summarizer = LsaSummarizer(stemmer)
        summarizer.stop_words = get_stop_words(LANGUAGE)
        summary = summarizer(parser.document, sentences_count)
        return " ".join(map(str, summary))
    except Exception as e:
        logging.warning(f"Summarization failed: {e}. Falling back to truncation.")
        # Fallback: return first few sentences if summarizer fails
        try:
            sentences = sent_tokenize(text) # Requires punkt tokenizer
            return " ".join(sentences[:sentences_count])
        except Exception: # If even tokenization fails
             return text[:500] + ("..." if len(text) > 500 else "")

def extract_keywords(text, max_keywords=MAX_KEYWORDS):
    if not text or not isinstance(text, str) or len(text.strip()) == 0:
        return ""
    try:
        words = word_tokenize(text.lower())
        # Improved filtering: longer than 2 chars, alpha, not stopword
        filtered_words = [word for word in words if len(word) > 2 and word.isalpha() and word not in stop_words_set]
        if not filtered_words:
            # logging.debug("No filtered words for keyword extraction.") # Optional debug
            return ""
        fdist = FreqDist(filtered_words)
        keywords = [word for word, freq in fdist.most_common(max_keywords)]
        # logging.debug(f"Extracted keywords: {keywords}") # Optional debug
        return ",".join(keywords)
    except LookupError as e:
        # This specific error happens if NLTK data (like 'punkt') isn't found
        log_msg = f"Keyword extraction failed due to missing NLTK data: {e}. Searched paths: {nltk.data.path}"
        print(f"\nERROR: {log_msg}", file=sys.stderr)
        logging.error(log_msg)
        # Ensure the check_nltk_data function ran or try downloading again.
        return "NLTK_DATA_ERROR"
    except Exception as e:
        log_msg = f"Keyword extraction failed: {e} (Type: {type(e).__name__})"
        # print(f"\nWarning: {log_msg}", file=sys.stderr) # Reduce console noise
        logging.warning(log_msg, exc_info=True) # Log with traceback
        return ""


# --- Database Interaction ---
def insert_update_db(cursor, file_data):
    """Inserts or updates a file record in the database."""
    sql = '''
        INSERT INTO files (
            path, filename, extension, size_bytes, last_modified,
            category_year, category_type, category_event, category_meeting,
            summary, keywords, processing_status, processing_error
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(path) DO UPDATE SET
            filename=excluded.filename,
            extension=excluded.extension,
            size_bytes=excluded.size_bytes,
            last_modified=excluded.last_modified,
            category_year=excluded.category_year,
            category_type=excluded.category_type,
            category_event=excluded.category_event,
            category_meeting=excluded.category_meeting,
            summary=excluded.summary,
            keywords=excluded.keywords,
            processing_status=excluded.processing_status,
            processing_error=excluded.processing_error
    '''
    try:
        cursor.execute(sql, (
            file_data.get('path'),
            file_data.get('filename'),
            file_data.get('extension'),
            file_data.get('size_bytes'),
            file_data.get('last_modified'),
            file_data.get('category_year'),
            file_data.get('category_type'),
            file_data.get('category_event', 'Unknown'),
            file_data.get('category_meeting', 'Unknown'),
            file_data.get('summary'),
            file_data.get('keywords'),
            file_data.get('processing_status', 'Failed'), # Default to Failed if not specified
            file_data.get('processing_error')
        ))
        return True
    except sqlite3.Error as e:
        log_msg = f"Database Error for '{file_data.get('path')}': {e}\nData: {file_data}"
        print(f"\n{log_msg}", file=sys.stderr)
        logging.error(log_msg)
        # Optionally add traceback: traceback.print_exc()
        return False

# --- Main Indexing Logic ---

def index_files(directory_path, db_conn, db_cursor):
    """Indexes files, extracting info and saving to DB."""
    if not os.path.isdir(directory_path):
        print(f"Error: Directory '{directory_path}' not found or is not accessible.", file=sys.stderr)
        logging.critical(f"Directory not found or inaccessible: {directory_path}")
        return 0, 0

    start_time = time.time()
    processed_count = 0
    failed_count = 0
    skipped_count = 0

    # --- Scan for all files first ---
    all_files_to_process = []
    print("Scanning directory tree...")
    try:
        for root, _, files in os.walk(directory_path, topdown=True, onerror=lambda e: logging.warning(f"Permission error walking '{e.filename}': {e}")):
            for filename in files:
                if filename.startswith('.'): # Skip hidden files
                    skipped_count += 1
                    continue
                full_path = os.path.join(root, filename)
                all_files_to_process.append(full_path)
    except Exception as e:
         print(f"Error during initial directory scan: {e}", file=sys.stderr)
         logging.error(f"Error during initial directory scan: {e}")
         # Might still be able to process the files found so far

    total_files = len(all_files_to_process)
    if total_files == 0 and skipped_count == 0:
        print("No processable files found.")
        return 0, 0

    print(f"Found {total_files} files to process (skipped {skipped_count} hidden files). Starting indexing...")

    # --- Process files with progress bar ---
    with tqdm(total=total_files, unit='file', desc="Indexing", mininterval=0.5, smoothing=0.1) as pbar:
        for i, file_path in enumerate(all_files_to_process):
            filename = os.path.basename(file_path)
            _, extension = os.path.splitext(filename)
            file_data = { # Initialize with path and defaults
                'path': file_path,
                'filename': filename,
                'extension': extension.lower() if extension else '',
                'processing_status': 'Failed', # Assume failure unless proven otherwise
                'processing_error': None,
                'summary': None,
                'keywords': None,
            }

            try:
                # --- Get Metadata ---
                stat_info = os.stat(file_path)
                mod_time = stat_info.st_mtime
                year = datetime.datetime.fromtimestamp(mod_time).year
                file_type = get_file_type(file_data['extension'])

                file_data.update({
                    'size_bytes': stat_info.st_size,
                    'last_modified': mod_time,
                    'category_year': year,
                    'category_type': file_type,
                })

                # --- Skip large files or specific types if needed ---
                # if stat_info.st_size > MAX_FILE_SIZE_MB * 1024 * 1024:
                #     file_data['processing_status'] = 'Skipped (Too Large)'
                #     skipped_count += 1
                #     insert_update_db(db_cursor, file_data)
                #     pbar.update(1)
                #     continue

                # --- Text Extraction ---
                extracted_text = None
                if file_type == 'Word Document':
                    extracted_text = extract_text_docx(file_path)
                elif file_type == 'Excel Spreadsheet':
                    extracted_text = extract_text_xlsx(file_path)
                elif file_type == 'PowerPoint Presentation':
                    extracted_text = extract_text_pptx(file_path)
                elif file_type == 'PDF Document':
                    extracted_text = extract_text_pdf(file_path)
                # --- Temporarily disable OCR to save memory ---
                # elif file_type == 'Image':
                #     extracted_text = extract_text_image(file_path) # OCR
                elif file_type == 'Text':
                    extracted_text = extract_text_plain(file_path)
                # Add more handlers here (e.g., 'Code', 'Archive' - maybe list contents?)

                # --- Summarization & Keywords ---
                if extracted_text:
                    file_data['summary'] = summarize_text(extracted_text)
                    file_data['keywords'] = extract_keywords(extracted_text)
                elif file_type not in ['Image', 'Other', 'Archive', 'Audio', 'Video', 'Code'] and extracted_text is None:
                    # Indicate if text *should* have been extracted but wasn't
                    file_data['processing_error'] = f"Failed to extract text ({file_type})"
                    logging.warning(f"Text extraction failed for expected type: {file_path}")

                # If we got here without major error, mark as success
                file_data['processing_status'] = 'Success'
                processed_count += 1

            except FileNotFoundError:
                file_data['processing_error'] = "File not found (moved/deleted during scan?)"
                logging.warning(f"File not found during processing: {file_path}")
                failed_count += 1
            except PermissionError:
                file_data['processing_error'] = "Permission denied"
                logging.warning(f"Permission denied for file: {file_path}")
                failed_count += 1
            except KeyboardInterrupt:
                 print("\nKeyboard interrupt detected. Committing progress and exiting.")
                 logging.warning("Keyboard interrupt detected.")
                 db_conn.commit() # Commit progress before exiting
                 raise # Re-raise to stop the program
            except Exception as e:
                file_data['processing_error'] = f"Unexpected error: {type(e).__name__}"
                log_msg = f"Error processing file '{file_path}': {e}"
                print(f"\n{log_msg}", file=sys.stderr)
                traceback.print_exc(limit=1, file=sys.stderr)
                logging.error(log_msg, exc_info=True)
                failed_count += 1
            finally:
                # Always attempt to record the file's status
                insert_update_db(db_cursor, file_data)
                pbar.update(1)

                # Commit periodically
                if (i + 1) % COMMIT_INTERVAL == 0:
                    db_conn.commit()
                    pbar.set_postfix_str("Committing...") # Show commit in progress bar

    # Final commit
    db_conn.commit()

    end_time = time.time()
    duration = end_time - start_time
    print(f"\n--- Indexing Summary ---")
    print(f"Total files scanned: {total_files + skipped_count}")
    print(f"Successfully processed: {processed_count}")
    print(f"Failed to process: {failed_count}")
    print(f"Skipped (hidden/other): {skipped_count}")
    print(f"Indexing completed in {duration:.2f} seconds ({duration/60:.2f} minutes).")
    print(f"Index data stored in '{DATABASE_NAME}'")
    print(f"Detailed errors logged in '{LOG_FILE}'")

    return processed_count, failed_count


# --- Main Execution ---
if __name__ == "__main__":
    print("Starting File Indexer...")
    if len(sys.argv) != 2:
        print(f"Usage: python {os.path.basename(__file__)} <directory_path>")
        sys.exit(1)

    target_directory = sys.argv[1]
    print(f"Target Directory: {target_directory}")
    print(f"Database File: {DATABASE_NAME}")
    print(f"Log File: {LOG_FILE}")

    # --- Prerequisite Checks ---
    if not check_nltk_data():
        sys.exit(1)
    # Add Tesseract check?
    if Image and pytesseract and getattr(pytesseract, 'pytesseract', None):
        try:
            pytesseract.get_tesseract_version()
            print(f"Tesseract Version: {pytesseract.get_tesseract_version()}")
        except (pytesseract.TesseractNotFoundError, Exception) as e:
             print("\nWarning: Tesseract OCR engine not found or accessible.", file=sys.stderr)
             print("Image file text extraction (OCR) will be skipped.", file=sys.stderr)
             print(f"(Error: {e})", file=sys.stderr)
             # Allow to continue without OCR


    # --- Database Setup & Indexing ---
    conn, cursor = None, None
    try:
        conn, cursor = setup_database()
        print(f"Database '{DATABASE_NAME}' initialized/opened.")

        # --- Start Indexing ---
        processed, failed = index_files(target_directory, conn, cursor)

    except sqlite3.Error as e:
        print(f"\nDatabase runtime error: {e}", file=sys.stderr)
        logging.critical(f"Database runtime error: {e}", exc_info=True)
    except KeyboardInterrupt:
         print("\nExiting due to user interrupt.")
         # Progress should have been committed in index_files or finally block
    except Exception as e:
        print(f"\nAn unexpected critical error occurred: {e}", file=sys.stderr)
        logging.critical(f"Unexpected critical error: {e}", exc_info=True)
        traceback.print_exc()
    finally:
        if conn:
            print("Committing final changes and closing database connection...")
            conn.commit()
            conn.close()
            print("Database connection closed.")

    print("Indexing process finished.")

